from pathlib import Path
import re

from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor

from pptx import Presentation
from pptx.dml.color import RGBColor as PPTColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches as PInches, Pt as PPt


ROOT = Path(__file__).resolve().parent.parent
REPORT_MD = ROOT / "phase1_report.md"
REPORT_DOCX = ROOT / "phase1_report_acl.docx"
SLIDES_PPTX = ROOT / "ManoVarta_Phase1_Presentation.pptx"


NAVY = PPTColor(28, 45, 72)
TEAL = PPTColor(42, 157, 143)
SAND = PPTColor(244, 241, 232)
GOLD = PPTColor(233, 196, 106)
CORAL = PPTColor(231, 111, 81)
SLATE = PPTColor(91, 109, 135)
WHITE = PPTColor(255, 255, 255)
TEXT = PPTColor(34, 34, 34)


def clean_inline_md(text: str) -> str:
    text = re.sub(r"\*\*(.*?)\*\*", r"\1", text)
    text = re.sub(r"`([^`]*)`", r"\1", text)
    text = re.sub(r"\[(.*?)\]\((.*?)\)", r"\1", text)
    return text.strip()


def set_section_columns(section, num=2, space_twips=360):
    sect_pr = section._sectPr
    cols = sect_pr.xpath("./w:cols")
    if cols:
        cols = cols[0]
    else:
        cols = OxmlElement("w:cols")
        sect_pr.append(cols)
    cols.set(qn("w:num"), str(num))
    cols.set(qn("w:space"), str(space_twips))


def set_doc_style(doc: Document):
    section = doc.sections[0]
    section.top_margin = Inches(1)
    section.bottom_margin = Inches(1)
    section.left_margin = Inches(1)
    section.right_margin = Inches(1)

    normal = doc.styles["Normal"]
    normal.font.name = "Times New Roman"
    normal.font.size = Pt(10)
    normal.paragraph_format.space_after = Pt(4)
    normal.paragraph_format.line_spacing = 1.0

    for style_name, size in [("Heading 1", 11), ("Heading 2", 10.5), ("Heading 3", 10)]:
        style = doc.styles[style_name]
        style.font.name = "Times New Roman"
        style.font.bold = True
        style.font.size = Pt(size)
        style.paragraph_format.space_before = Pt(6)
        style.paragraph_format.space_after = Pt(3)


def add_centered(doc: Document, text: str, size: int, bold=False, italic=False):
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(text)
    run.font.name = "Times New Roman"
    run.font.size = Pt(size)
    run.bold = bold
    run.italic = italic
    return p


def add_paragraph(doc: Document, text: str, justify=True):
    p = doc.add_paragraph()
    if justify:
        p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
    run = p.add_run(clean_inline_md(text))
    run.font.name = "Times New Roman"
    run.font.size = Pt(10)
    return p


def add_code_block(doc: Document, code_lines):
    p = doc.add_paragraph()
    p.paragraph_format.left_indent = Inches(0.2)
    p.paragraph_format.space_before = Pt(3)
    p.paragraph_format.space_after = Pt(3)
    for idx, line in enumerate(code_lines):
        run = p.add_run(line)
        if idx != len(code_lines) - 1:
            run.add_break()
        run.font.name = "Courier New"
        run.font.size = Pt(8.5)


def add_md_table(doc: Document, table_lines):
    rows = []
    for line in table_lines:
        if re.match(r"^\|\s*-", line):
            continue
        cells = [clean_inline_md(cell) for cell in line.strip().strip("|").split("|")]
        rows.append([c.strip() for c in cells])
    if not rows:
        return
    table = doc.add_table(rows=len(rows), cols=len(rows[0]))
    table.style = "Table Grid"
    table.autofit = True
    for r_idx, row in enumerate(rows):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = value
            for paragraph in cell.paragraphs:
                paragraph.paragraph_format.space_after = Pt(0)
                if r_idx == 0:
                    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                for run in paragraph.runs:
                    run.font.name = "Times New Roman"
                    run.font.size = Pt(9)
                    if r_idx == 0:
                        run.bold = True
    doc.add_paragraph()


def build_docx():
    doc = Document()
    set_doc_style(doc)

    add_centered(doc, "ManoVarta: Multilingual Conversational AI Chatbot for Mental Health Screening", 16, bold=True)
    add_centered(doc, "Ritwij and Yash", 11, bold=True)
    add_centered(doc, "Graduate NLP/AI Project", 10)
    add_centered(doc, "ACL-style Phase 1 Report Draft", 10, italic=True)
    doc.add_paragraph()

    lines = REPORT_MD.read_text().splitlines()
    try:
        start_idx = lines.index("## Abstract")
    except ValueError as exc:
        raise RuntimeError("Could not locate abstract heading in phase1_report.md") from exc

    # Parse front matter and main content.
    main_section_created = False
    paragraph_buffer = []
    i = start_idx

    def flush_paragraph():
        nonlocal paragraph_buffer
        if paragraph_buffer:
            add_paragraph(doc, " ".join(line.strip() for line in paragraph_buffer))
            paragraph_buffer = []

    while i < len(lines):
        line = lines[i].rstrip()
        stripped = line.strip()

        if stripped == "":
            flush_paragraph()
            i += 1
            continue

        if stripped == "## 1. Introduction" and not main_section_created:
            flush_paragraph()
            new_section = doc.add_section(WD_SECTION.CONTINUOUS)
            set_section_columns(new_section, num=2, space_twips=360)
            main_section_created = True

        if stripped.startswith("## "):
            flush_paragraph()
            heading = clean_inline_md(stripped[3:])
            p = doc.add_paragraph(style="Heading 1")
            p.alignment = WD_ALIGN_PARAGRAPH.LEFT
            run = p.add_run(heading)
            run.font.name = "Times New Roman"
            run.font.size = Pt(11)
            run.bold = True
            i += 1
            continue

        if stripped.startswith("### "):
            flush_paragraph()
            heading = clean_inline_md(stripped[4:])
            p = doc.add_paragraph(style="Heading 2")
            p.alignment = WD_ALIGN_PARAGRAPH.LEFT
            run = p.add_run(heading)
            run.font.name = "Times New Roman"
            run.font.size = Pt(10.5)
            run.bold = True
            i += 1
            continue

        if stripped.startswith("```"):
            flush_paragraph()
            code_lines = []
            i += 1
            while i < len(lines) and not lines[i].strip().startswith("```"):
                code_lines.append(lines[i])
                i += 1
            add_code_block(doc, code_lines)
            i += 1
            continue

        if stripped.startswith("|"):
            flush_paragraph()
            table_lines = []
            while i < len(lines) and lines[i].strip().startswith("|"):
                table_lines.append(lines[i].rstrip())
                i += 1
            add_md_table(doc, table_lines)
            continue

        if stripped.startswith("- "):
            flush_paragraph()
            while i < len(lines) and lines[i].strip().startswith("- "):
                bullet = clean_inline_md(lines[i].strip()[2:])
                p = doc.add_paragraph(style="List Bullet")
                p.paragraph_format.space_after = Pt(1)
                run = p.add_run(bullet)
                run.font.name = "Times New Roman"
                run.font.size = Pt(10)
                i += 1
            continue

        if re.match(r"^\d+\.\s", stripped):
            flush_paragraph()
            while i < len(lines) and re.match(r"^\d+\.\s", lines[i].strip()):
                item = clean_inline_md(re.sub(r"^\d+\.\s", "", lines[i].strip()))
                p = doc.add_paragraph(style="List Number")
                p.paragraph_format.space_after = Pt(1)
                run = p.add_run(item)
                run.font.name = "Times New Roman"
                run.font.size = Pt(10)
                i += 1
            continue

        paragraph_buffer.append(stripped)
        i += 1

    flush_paragraph()
    doc.save(REPORT_DOCX)


def add_slide_base(slide, title, subtitle=None, idx=None):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = SAND

    header = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, PInches(0), PInches(0), PInches(13.333), PInches(0.55)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = NAVY
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(PInches(0.45), PInches(0.18), PInches(8.8), PInches(0.3))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Calibri"
    p.font.size = PPt(27)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if subtitle:
        sub_box = slide.shapes.add_textbox(PInches(9.25), PInches(0.17), PInches(3.7), PInches(0.28))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        p.text = subtitle
        p.font.name = "Calibri"
        p.font.size = PPt(12)
        p.font.color.rgb = WHITE

    if idx is not None:
        footer = slide.shapes.add_textbox(PInches(12.8), PInches(7.05), PInches(0.3), PInches(0.2))
        tf = footer.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        p.text = str(idx)
        p.font.name = "Calibri"
        p.font.size = PPt(10)
        p.font.color.rgb = SLATE


def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = text
    p.font.name = "Calibri"
    p.font.size = PPt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, height, bullets, font_size=19, color=TEXT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.clear()
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.bullet = True
        p.font.name = "Calibri"
        p.font.size = PPt(font_size)
        p.font.color.rgb = color
    return box


def add_card(slide, left, top, width, height, title, body, fill=WHITE, title_color=NAVY, body_color=TEXT):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = TEAL
    shape.line.width = PPt(1.5)

    add_textbox(slide, left + PInches(0.15), top + PInches(0.12), width - PInches(0.3), PInches(0.28), title, font_size=17, bold=True, color=title_color)
    add_textbox(slide, left + PInches(0.15), top + PInches(0.42), width - PInches(0.3), height - PInches(0.5), body, font_size=13, color=body_color)
    return shape


def add_table(slide, left, top, width, height, rows, header_fill=NAVY):
    cols = len(rows[0])
    table_shape = slide.shapes.add_table(len(rows), cols, left, top, width, height)
    table = table_shape.table
    for r_idx, row in enumerate(rows):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = value
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if r_idx == 0 else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.name = "Calibri"
                    run.font.size = PPt(12 if r_idx == 0 else 11)
                    run.font.bold = r_idx == 0
                    run.font.color.rgb = WHITE if r_idx == 0 else TEXT
    return table


def build_pptx():
    prs = Presentation()
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_base(slide, "ManoVarta", "Phase 1 Proposal", idx=1)
    add_textbox(slide, PInches(0.65), PInches(1.05), PInches(9.0), PInches(0.8),
                "Multilingual Conversational AI for Mental Health Screening", font_size=28, bold=True, color=NAVY)
    add_textbox(slide, PInches(0.68), PInches(1.9), PInches(6.3), PInches(0.6),
                "English + Hindi | Text-first | Screening support for professionals", font_size=18, color=SLATE)
    for idx, (title, body, left) in enumerate([
        ("Conversation", "Open-ended but bounded interaction", 0.8),
        ("Evidence", "Traceable PHQ-9 and GAD-7 item support", 4.55),
        ("Summary", "Clinician-facing scores, confidence, and safety flags", 8.3),
    ]):
        add_card(slide, PInches(left), PInches(3.2), PInches(3.0), PInches(1.7), title, body, fill=WHITE)
        if idx < 2:
            arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, PInches(left + 3.05), PInches(3.65), PInches(0.5), PInches(0.6))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = GOLD
            arrow.line.fill.background()
    add_textbox(slide, PInches(0.8), PInches(6.1), PInches(11.5), PInches(0.5),
                "Phase 1 goal: define the research problem, data plan, architecture, and evaluation strategy. Not a deployed clinical product.", font_size=15, color=SLATE)

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_base(slide, "Why This Problem", idx=2)
    add_table(slide, PInches(0.75), PInches(1.15), PInches(5.8), PInches(2.6), [
        ["Direct form problem", "Why conversation may help"],
        ["survey fatigue", "more natural disclosure"],
        ["guarded answers", "room for follow-up"],
        ["rigid item order", "adaptive symptom coverage"],
    ])
    add_card(slide, PInches(7.2), PInches(1.3), PInches(5.0), PInches(2.25), "Core tension",
             "Mental health screening needs clinical structure, but users often describe symptoms in indirect, messy, multilingual ways.")
    add_bullets(slide, PInches(7.25), PInches(3.95), PInches(5.0), PInches(1.6), [
        "Questionnaires are useful, but can feel impersonal.",
        "Natural conversation may improve disclosure.",
        "The challenge is preserving traceable item-level structure."
    ], font_size=16)

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_base(slide, "Objective and Scope", idx=3)
    add_table(slide, PInches(0.8), PInches(1.25), PInches(6.0), PInches(2.6), [
        ["In scope", "Out of scope for Phase 1"],
        ["English + Hindi text chat", "voice assistant"],
        ["PHQ-9 and GAD-7 grounding", "therapy or diagnosis"],
        ["item-level score inference", "production deployment"],
    ], header_fill=TEAL)
    add_card(slide, PInches(7.3), PInches(1.35), PInches(4.9), PInches(2.2), "Main objective",
             "Infer symptom evidence and questionnaire-aligned item scores from natural conversation while keeping human oversight in the loop.")
    add_textbox(slide, PInches(7.35), PInches(4.0), PInches(5.0), PInches(1.4),
                "Phase 1 is a foundation milestone: proposal, data design, architecture, baselines, and evaluation plan.", font_size=18, color=SLATE)

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_base(slide, "Research Questions", idx=4)
    questions = [
        "Can natural dialogue recover PHQ-9 and GAD-7 item scores?",
        "Does confidence-based follow-up improve coverage?",
        "Is performance stable across English, Hindi, and Hinglish?",
    ]
    for i, q in enumerate(questions):
        add_card(slide, PInches(0.9 + i * 4.05), PInches(1.6), PInches(3.6), PInches(2.0),
                 f"RQ{i+1}", q, fill=WHITE)
    add_card(slide, PInches(2.0), PInches(4.3), PInches(9.2), PInches(1.65), "Contribution",
             "A modular bilingual screening setup with rapport-aware dialogue, evidence-first scoring, item confidence tracking, and a parallel safety module.")

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_base(slide, "Example: Conversation to Structured Output", idx=5)
    bubble = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, PInches(0.8), PInches(1.45), PInches(5.2), PInches(1.3))
    bubble.fill.solid()
    bubble.fill.fore_color.rgb = WHITE
    bubble.line.color.rgb = TEAL
    add_textbox(slide, PInches(1.0), PInches(1.7), PInches(4.8), PInches(0.8),
                "User: \"My sleep schedule is messed up and I can't focus on assignments.\"", font_size=18, color=NAVY)
    arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, PInches(6.2), PInches(1.8), PInches(0.7), PInches(0.6))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = GOLD
    arrow.line.fill.background()
    add_table(slide, PInches(7.1), PInches(1.35), PInches(5.2), PInches(2.2), [
        ["Evidence span", "Item", "Score"],
        ["sleep schedule is messed up", "PHQ-9 sleep", "2"],
        ["can't focus on assignments", "PHQ-9 concentration", "2"],
    ], header_fill=TEAL)
    add_card(slide, PInches(1.0), PInches(4.25), PInches(4.5), PInches(1.5), "Confidence state",
             "sleep: high\nconcentration: medium-high\nlow mood: unresolved")
    add_card(slide, PInches(7.0), PInches(4.25), PInches(4.6), PInches(1.5), "Next action",
             "Ask a focused follow-up about mood or fatigue instead of repeating the whole questionnaire.")

    # Slide 6
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_base(slide, "Data Plan", idx=6)
    add_table(slide, PInches(0.8), PInches(1.25), PInches(4.4), PInches(2.0), [
        ["Subset", "Conversations"],
        ["English", "32"],
        ["Hindi", "32"],
        ["Hinglish", "16"],
    ])
    chart_left = 6.0
    max_width = 5.0
    for i, (label, value, color) in enumerate([("English", 32, TEAL), ("Hindi", 32, NAVY), ("Hinglish", 16, GOLD)]):
        y = 1.4 + i * 0.8
        add_textbox(slide, PInches(chart_left), PInches(y), PInches(1.1), PInches(0.25), label, font_size=14, color=NAVY)
        bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, PInches(chart_left + 1.15), PInches(y + 0.05), PInches(max_width * value / 32), PInches(0.28))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        add_textbox(slide, PInches(chart_left + 1.25 + max_width * value / 32), PInches(y), PInches(0.6), PInches(0.25), str(value), font_size=14, color=SLATE)
    add_card(slide, PInches(0.9), PInches(4.0), PInches(11.3), PInches(1.45), "Pilot data policy",
             "40 synthetic patient profiles -> 80 total conversations. Public datasets are used only for auxiliary validation and safety robustness, not as a substitute for item-level PHQ-9/GAD-7 annotation.")

    # Slide 7
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_base(slide, "Annotation Workflow", idx=7)
    steps = [
        ("1. Profile design", "demographics, stressors, symptom pattern"),
        ("2. Dialogue drafting", "guided generation + manual revision"),
        ("3. Double annotation", "PHQ/GAD labels, spans, safety"),
        ("4. Consensus review", "resolve disagreements and freeze gold labels"),
    ]
    for i, (title, body) in enumerate(steps):
        left = 0.65 + i * 3.15
        add_card(slide, PInches(left), PInches(2.0), PInches(2.6), PInches(2.0), title, body)
        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, PInches(left + 2.65), PInches(2.65), PInches(0.35), PInches(0.45))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = GOLD
            arrow.line.fill.background()
    add_textbox(slide, PInches(0.85), PInches(5.0), PInches(11.5), PInches(0.6),
                "Quality checks: weighted kappa for item labels, span agreement, and explicit notes on ambiguity or contradiction.", font_size=17, color=SLATE)

    # Slide 8
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_base(slide, "System Architecture", idx=8)
    add_card(slide, PInches(0.9), PInches(2.1), PInches(3.2), PInches(1.8), "Rapport-aware dialogue manager",
             "Handles greeting, topic steering, and adaptive follow-up strategy.")
    add_card(slide, PInches(5.0), PInches(2.1), PInches(3.3), PInches(1.8), "Evidence and scoring engine",
             "Extracts spans, maps to PHQ-9/GAD-7 items, assigns scores and confidence.")
    add_card(slide, PInches(9.2), PInches(1.5), PInches(3.0), PInches(2.1), "Parallel safety trigger",
             "Flags crisis-sensitive language for review independently of item scoring.", fill=WHITE, title_color=CORAL)
    for x in [4.2, 8.45]:
        arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, PInches(x), PInches(2.65), PInches(0.45), PInches(0.55))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = GOLD
        arrow.line.fill.background()
    add_card(slide, PInches(3.9), PInches(4.55), PInches(5.2), PInches(1.35), "State tracking",
             "LangGraph-style orchestration keeps coverage, confidence, and escalation decisions explicit.")

    # Slide 9
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_base(slide, "Model Stack", idx=9)
    add_table(slide, PInches(0.85), PInches(1.25), PInches(8.0), PInches(3.2), [
        ["Component", "Model"],
        ["main dialogue + evidence", "Aya Expanse 32B"],
        ["open comparison", "Mistral NeMo 12B"],
        ["optional second baseline", "Gemma 3 12B"],
        ["safety encoder", "IndicBERT-style model"],
    ], header_fill=TEAL)
    add_card(slide, PInches(9.2), PInches(1.35), PInches(3.1), PInches(2.8), "Why this stack?",
             "Aya is the primary multilingual model.\nMistral NeMo is the practical open comparison.\nThe encoder keeps Hindi-sensitive safety separate.")
    add_textbox(slide, PInches(0.9), PInches(5.2), PInches(11.0), PInches(0.6),
                "Phase 1 does not assume full fine-tuning of a 32B model. Most iteration can happen on smaller open baselines.", font_size=16, color=SLATE)

    # Slide 10
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_base(slide, "Evaluation Plan", idx=10)
    metric_cards = [
        ("Item-level MAE / Macro-F1", "Do predicted item scores match the gold labels?"),
        ("Evidence support rate", "Are score predictions backed by valid evidence spans?"),
        ("Safety recall / precision", "Does the system catch crisis-sensitive language?"),
        ("Disclosure efficiency", "How quickly do uncertain items become resolved?"),
    ]
    positions = [(0.9, 1.4), (6.9, 1.4), (0.9, 3.85), (6.9, 3.85)]
    for (title, body), (x, y) in zip(metric_cards, positions):
        add_card(slide, PInches(x), PInches(y), PInches(5.3), PInches(1.75), title, body)
    add_textbox(slide, PInches(0.95), PInches(6.2), PInches(11.0), PInches(0.45),
                "Also tracked: coverage completeness, multilingual parity, and latency as a future integration metric.", font_size=15, color=SLATE)

    # Slide 11
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_base(slide, "Baselines and Ablations", idx=11)
    add_table(slide, PInches(0.8), PInches(1.3), PInches(11.8), PInches(3.0), [
        ["Comparison set", "Adaptive?", "Evidence-first?", "Safety-aware?", "Purpose"],
        ["Direct questionnaire", "No", "No", "Manual only", "strongest structured baseline"],
        ["Fixed scripted chatbot", "No", "Partial", "Optional", "conversation without adaptivity"],
        ["Single-pass transcript scoring", "No", "No", "No", "cheap transcript-only scorer"],
        ["No confidence / no safety", "Partial", "Yes", "No", "ablation checks"],
    ], header_fill=NAVY)
    add_textbox(slide, PInches(0.95), PInches(5.0), PInches(11.0), PInches(0.8),
                "Main claim to test: the adaptive evidence-first system should improve coverage and interpretability without sacrificing item-level scoring quality.", font_size=16, color=SLATE)

    # Slide 12
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_base(slide, "Risks and Ethics", idx=12)
    add_table(slide, PInches(0.85), PInches(1.25), PInches(11.7), PInches(3.55), [
        ["Risk", "Mitigation"],
        ["small pilot data", "keep claims modest and emphasize pilot nature"],
        ["synthetic vs real-world gap", "manual review and explicit realism limits"],
        ["Hindi nuance and code-mixing", "track parity and include code-mixed cases early"],
        ["safety misses", "parallel safety module + human review"],
    ], header_fill=CORAL)
    add_textbox(slide, PInches(0.95), PInches(5.2), PInches(11.2), PInches(0.9),
                "The system is a screening support tool, not a diagnostic or therapeutic replacement. Any later real-user study would require consent, privacy safeguards, and institutional review where appropriate.", font_size=16, color=SLATE)

    # Slide 13
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_base(slide, "Milestones and Team Split", idx=13)
    timeline = [
        ("Phase 1", "proposal + data/eval design", TEAL, 0.9),
        ("Phase 2", "prototype pipeline + pilot expansion", NAVY, 4.5),
        ("Phase 3", "experiments + demo + final report", GOLD, 8.5),
    ]
    for title, body, color, left in timeline:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, PInches(left), PInches(1.45), PInches(3.1), PInches(1.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        add_textbox(slide, PInches(left + 0.15), PInches(1.65), PInches(2.8), PInches(0.25), title, font_size=18, bold=True, color=WHITE)
        add_textbox(slide, PInches(left + 0.15), PInches(2.05), PInches(2.8), PInches(0.5), body, font_size=13, color=WHITE)
    add_card(slide, PInches(1.6), PInches(4.15), PInches(4.4), PInches(1.6), "Ritwij",
             "inference engine\nscoring logic\nevaluation and integration")
    add_card(slide, PInches(7.1), PInches(4.15), PInches(4.4), PInches(1.6), "Yash",
             "data design\nannotation workflow\ndialogue and architecture")
    add_textbox(slide, PInches(1.0), PInches(6.2), PInches(11.0), PInches(0.4),
                "Immediate next step: freeze the annotation rubric and build the first pilot batch of conversations.", font_size=16, color=SLATE, align=PP_ALIGN.CENTER)

    prs.save(SLIDES_PPTX)


if __name__ == "__main__":
    build_docx()
    build_pptx()
    print(f"Wrote {REPORT_DOCX.name}")
    print(f"Wrote {SLIDES_PPTX.name}")
