from pathlib import Path
import tempfile
import zipfile

from pptx import Presentation
from pptx.chart.data import ChartData, CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION, XL_LEGEND_POSITION, XL_MARKER_STYLE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path("/Users/ritwij/Documents/multilingualChatbot")
TEMPLATE_ZIP = Path("/Users/ritwij/Documents/Formal_Powerpoint_Master_Brand_Widescreen.zip")
OUTPUT = ROOT / "ManoVarta_Phase1_Presentation_UBTemplate.pptx"

UB_BLUE = RGBColor(0x00, 0x5B, 0xBB)
UB_LT_BLUE = RGBColor(0x2F, 0x9F, 0xD0)
UB_ORANGE = RGBColor(0xE5, 0x69, 0x54)
UB_DK_BLUE = RGBColor(0x00, 0x2F, 0x56)
UB_TEAL = RGBColor(0x00, 0x65, 0x70)
UB_GRAY = RGBColor(0x4D, 0x4D, 0x4D)
UB_LT_GRAY = RGBColor(0xE3, 0xE3, 0xE3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


LAYOUT_TITLE = 0
LAYOUT_DIVIDER = 1
LAYOUT_ONE_COL_BULLETS = 2
LAYOUT_ONE_COL_TEXT = 4
LAYOUT_TWO_COL_TEXT = 5
LAYOUT_TWO_COL_BULLETS = 6
LAYOUT_TWO_COL_SUB_BULLETS = 8
LAYOUT_TEXT_CHART = 9
LAYOUT_BLANK = 10


def extract_template() -> Path:
    with zipfile.ZipFile(TEMPLATE_ZIP) as zf:
        name = next(n for n in zf.namelist() if n.endswith(".pptx") and "__MACOSX" not in n)
        tmpdir = Path(tempfile.mkdtemp(prefix="ub_template_"))
        zf.extract(name, tmpdir)
        return tmpdir / name


def remove_slide(prs: Presentation, index: int):
    slide_id = prs.slides._sldIdLst[index]
    r_id = slide_id.rId
    prs.part.drop_rel(r_id)
    prs.slides._sldIdLst.remove(slide_id)


def clear_all_template_slides(prs: Presentation):
    for idx in reversed(range(len(prs.slides))):
        remove_slide(prs, idx)


def set_text(shape, text, font_size=None, bold=None, color=None, align=None, font_name="Arial"):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if align is not None:
        p.alignment = align
    if p.runs:
        run = p.runs[0]
        run.font.name = font_name
        if font_size is not None:
            run.font.size = Pt(font_size)
        if bold is not None:
            run.font.bold = bold
        if color is not None:
            run.font.color.rgb = color
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE


def set_bullets(shape, bullets, font_size=24, color=UB_GRAY, levels=None):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0 if levels is None else levels[i]
        if p.runs:
            run = p.runs[0]
            run.font.name = "Arial"
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE


def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=UB_GRAY, align=PP_ALIGN.LEFT, font_name="Arial"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    if p.runs:
        run = p.runs[0]
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_multiline_box(slide, left, top, width, height, lines, font_size=16, color=UB_GRAY, font_name="Arial", bold_first=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        if p.runs:
            run = p.runs[0]
            run.font.name = font_name
            run.font.size = Pt(font_size)
            run.font.bold = bold_first and i == 0
            run.font.color.rgb = color
    return box


def add_card(slide, left, top, width, height, title, body, fill=WHITE, line=UB_BLUE, title_color=UB_DK_BLUE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.25)
    add_textbox(slide, left + Inches(0.16), top + Inches(0.10), width - Inches(0.32), Inches(0.28), title, font_size=16, bold=True, color=title_color)
    add_textbox(slide, left + Inches(0.16), top + Inches(0.42), width - Inches(0.32), height - Inches(0.50), body, font_size=13, color=UB_GRAY)
    return shape


def add_table(slide, left, top, width, height, data, header_fill=UB_BLUE):
    rows, cols = len(data), len(data[0])
    shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = shape.table
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c]
            cell.text_frame.word_wrap = True
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_fill if r == 0 else WHITE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if r == 0 else PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.name = "Arial"
                    run.font.size = Pt(12 if r == 0 else 11)
                    run.font.bold = r == 0
                    run.font.color.rgb = WHITE if r == 0 else UB_GRAY
    return table


def get_content_box(slide):
    ph = slide.placeholders[1]
    return ph.left, ph.top, ph.width, ph.height


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE])
    set_text(slide.shapes.title, "ManoVarta", font_size=42, bold=True, color=WHITE)
    subtitle = slide.placeholders[1]
    tf = subtitle.text_frame
    tf.clear()
    lines = [
        "Multilingual Conversational AI for Mental Health Screening",
        "Phase 1 proposal | English + Hindi | Ritwij and Yash",
    ]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        if p.runs:
            run = p.runs[0]
            run.font.name = "Arial"
            run.font.size = Pt(22 if i == 0 else 18)
            run.font.color.rgb = WHITE
    return slide


def add_problem_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TWO_COL_BULLETS])
    slide.shapes.title.text = "Why This Problem Matters"
    set_bullets(slide.placeholders[1], [
        "direct forms can create survey fatigue",
        "users may give guarded or minimal answers",
        "fixed item order misses natural symptom narratives",
    ], font_size=23)
    set_bullets(slide.placeholders[13], [
        "conversation can invite more natural disclosure",
        "follow-up can target missing or ambiguous symptoms",
        "the challenge is keeping clinical structure intact",
    ], font_size=23)
    left, top, width, height = get_content_box(slide)
    add_textbox(slide, Inches(0.55), Inches(6.15), Inches(12.0), Inches(0.35),
                "Questionnaires remain important. The research question is whether conversational elicitation can preserve PHQ-9/GAD-7 structure while sounding less rigid.",
                font_size=15, color=UB_DK_BLUE, bold=False)
    return slide


def add_gap_questions_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TWO_COL_SUB_BULLETS])
    slide.shapes.title.text = "Research Gap and Phase 1 Questions"
    set_text(slide.placeholders[1], "Known gaps", font_size=18, bold=True, color=UB_BLUE)
    set_bullets(slide.placeholders[14], [
        "small pilot data in realistic bilingual settings",
        "public labels rarely map cleanly to item-level PHQ-9/GAD-7 scoring",
        "Hindi nuance and code-mixing are often treated as afterthoughts",
    ], font_size=22)
    set_text(slide.placeholders[15], "Phase 1 questions", font_size=18, bold=True, color=UB_BLUE)
    set_bullets(slide.placeholders[16], [
        "Can natural dialogue recover item-level symptom scores?",
        "Does confidence-based follow-up improve coverage?",
        "How stable is performance across English, Hindi, and Hinglish?",
    ], font_size=22)
    return slide


def add_example_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_ONE_COL_BULLETS])
    slide.shapes.title.text = "Example: Conversation to Structured Output"
    slide.placeholders[1].text = ""
    left, top, width, height = get_content_box(slide)

    add_card(slide, Inches(0.60), Inches(2.05), Inches(3.0), Inches(2.8), "Sample patient profile",
             "23-year-old graduate student\nEnglish-dominant\nmoderate sleep disruption\nconcentration difficulty\nno immediate safety flag")

    bubble = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(4.05), Inches(2.0), Inches(4.1), Inches(1.15))
    bubble.fill.solid()
    bubble.fill.fore_color.rgb = WHITE
    bubble.line.color.rgb = UB_LT_BLUE
    bubble.line.width = Pt(1.5)
    add_textbox(slide, Inches(4.25), Inches(2.22), Inches(3.7), Inches(0.65),
                "User: \"My sleep schedule is messed up and I can't focus on assignments.\"", font_size=17, color=UB_DK_BLUE)

    arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(8.35), Inches(2.28), Inches(0.55), Inches(0.42))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = UB_ORANGE
    arrow.line.fill.background()

    add_table(slide, Inches(9.05), Inches(1.95), Inches(3.55), Inches(2.0), [
        ["Evidence span", "Item", "Score"],
        ["sleep schedule...", "PHQ-9 sleep", "2"],
        ["can't focus...", "PHQ-9 concentration", "2"],
    ], header_fill=UB_TEAL)

    add_card(slide, Inches(4.15), Inches(3.65), Inches(2.55), Inches(1.3), "Confidence state",
             "sleep: high\nconcentration: medium-high\nlow mood: unresolved", fill=WHITE, line=UB_TEAL)
    add_card(slide, Inches(7.10), Inches(3.65), Inches(4.7), Inches(1.3), "Next action",
             "Ask a targeted follow-up about mood or fatigue instead of repeating the full questionnaire.",
             fill=WHITE, line=UB_ORANGE)
    return slide


def add_architecture_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_ONE_COL_BULLETS])
    slide.shapes.title.text = "Proposed System Architecture"
    slide.placeholders[1].text = ""
    add_card(slide, Inches(0.70), Inches(2.1), Inches(2.6), Inches(1.55), "User turn",
             "English, Hindi, or Hinglish input")
    add_card(slide, Inches(3.65), Inches(2.1), Inches(2.9), Inches(1.55), "Rapport-aware dialogue manager",
             "topic steering\ncoverage tracking\nfollow-up strategy", line=UB_BLUE)
    add_card(slide, Inches(6.90), Inches(2.1), Inches(3.0), Inches(1.55), "Evidence and scoring engine",
             "span extraction\nitem mapping\n0-3 scoring + confidence", line=UB_TEAL)
    add_card(slide, Inches(10.25), Inches(2.1), Inches(2.35), Inches(1.55), "Clinician summary",
             "item scores\nconfidence\nsafety flags", line=UB_DK_BLUE)

    for x in [3.30, 6.55, 9.90]:
        arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(x), Inches(2.55), Inches(0.32), Inches(0.35))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = UB_ORANGE
        arrow.line.fill.background()

    add_card(slide, Inches(4.2), Inches(4.25), Inches(4.3), Inches(1.45), "Parallel safety trigger module",
             "runs independently of symptom scoring\nflags escalation-sensitive language for human review",
             fill=WHITE, line=UB_ORANGE, title_color=UB_ORANGE)
    dashed1 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(5.7), Inches(3.6), Inches(0.3), Inches(0.35))
    dashed1.fill.solid()
    dashed1.fill.fore_color.rgb = UB_LT_BLUE
    dashed1.line.fill.background()
    dashed2 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(8.0), Inches(3.6), Inches(0.3), Inches(0.35))
    dashed2.fill.solid()
    dashed2.fill.fore_color.rgb = UB_LT_BLUE
    dashed2.line.fill.background()
    add_textbox(slide, Inches(0.85), Inches(6.05), Inches(11.6), Inches(0.35),
                "Key design choice: safety runs in parallel rather than waiting for the main dialogue model to be correct.", font_size=15, color=UB_DK_BLUE)
    return slide


def add_data_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TEXT_CHART])
    slide.shapes.title.text = "Pilot Data Plan"
    set_bullets(slide.placeholders[2], [
        "40 synthetic patient profiles expanded into 80 conversations",
        "32 English, 32 Hindi, 16 Hinglish/code-mixed",
        "profile-guided drafting plus manual review for realism",
        "DAIC-WOZ, CLPsych, and eRisk only for auxiliary validation and safety stress tests",
    ], font_size=20)

    chart_data = ChartData()
    chart_data.categories = ["English", "Hindi", "Hinglish"]
    chart_data.add_series("Conversations", (32, 32, 16))
    chart = slide.placeholders[10].insert_chart(XL_CHART_TYPE.DOUGHNUT, chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    plot = chart.plots[0]
    plot.doughnut_hole_size = 58
    plot.has_data_labels = True
    labels = plot.data_labels
    labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
    labels.show_percentage = True
    labels.show_category_name = True
    labels.number_format = '0%'
    series = chart.series[0]
    for point, color in zip(series.points, [UB_BLUE, UB_LT_BLUE, UB_ORANGE]):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = color

    add_textbox(slide, Inches(0.65), Inches(6.05), Inches(11.8), Inches(0.3),
                "Planned pilot composition only; not presented as a completed clinical dataset.", font_size=14, color=UB_DK_BLUE)
    return slide


def add_annotation_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_ONE_COL_BULLETS])
    slide.shapes.title.text = "Annotation Protocol and Data Schema"
    slide.placeholders[1].text = ""

    steps = [
        ("1. Profile design", "age, context, symptom pattern, disclosure style"),
        ("2. Dialogue drafting", "guided generation followed by manual editing"),
        ("3. Double annotation", "PHQ-9, GAD-7, evidence spans, safety flag"),
        ("4. Consensus pass", "resolve disagreements and freeze gold labels"),
    ]
    for i, (title, body) in enumerate(steps):
        x = 0.6 + i * 3.05
        add_card(slide, Inches(x), Inches(1.95), Inches(2.6), Inches(1.65), title, body, line=UB_LT_BLUE if i % 2 else UB_BLUE)

    code_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(4.1), Inches(12.0), Inches(1.9))
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = WHITE
    code_box.line.color.rgb = UB_LT_GRAY
    add_multiline_box(slide, Inches(1.0), Inches(4.35), Inches(11.4), Inches(1.3), [
        '{',
        '  "patient_id": "MB-P001",   "language": "en",   "evidence_spans": ["sleep schedule is messed up"],',
        '  "phq_q3_sleep": 2,   "phq_q7_concentration": 2,   "safety_flag": "none"',
        '}',
    ], font_size=14, color=UB_DK_BLUE, font_name="Courier New")
    add_textbox(slide, Inches(0.85), Inches(6.2), Inches(11.5), Inches(0.28),
                "Agreement targets: weighted kappa for item labels, span overlap for evidence, and explicit notes on ambiguity or contradiction.",
                font_size=14, color=UB_DK_BLUE)
    return slide


def add_model_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_ONE_COL_BULLETS])
    slide.shapes.title.text = "Model Stack and Compute Rationale"
    slide.placeholders[1].text = ""

    blocks = [
        ("Primary dialogue + evidence", "Aya Expanse 32B", UB_BLUE),
        ("Open comparison model", "Mistral NeMo 12B", UB_LT_BLUE),
        ("Optional second baseline", "Gemma 3 12B", UB_TEAL),
        ("Safety / retrieval encoder", "IndicBERT-style model", UB_ORANGE),
    ]
    for i, (title, body, color) in enumerate(blocks):
        add_card(slide, Inches(0.8 + i * 3.05), Inches(2.0), Inches(2.6), Inches(1.6), title, body, line=color)

    add_card(slide, Inches(1.2), Inches(4.25), Inches(4.7), Inches(1.4), "Why modular?",
             "Dialogue, scoring, and safety should be separable so the system is easier to inspect, compare, and debug.", line=UB_TEAL)
    add_card(slide, Inches(6.45), Inches(4.25), Inches(5.1), Inches(1.4), "Compute plan",
             "Use smaller baselines for most iteration. Reserve Aya Expanse for targeted evaluation subsets, quantized runs, or hosted inference.", line=UB_ORANGE)
    return slide


def add_evaluation_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TEXT_CHART])
    slide.shapes.title.text = "Evaluation Plan"
    set_bullets(slide.placeholders[2], [
        "item-level MAE and macro-F1 for PHQ-9/GAD-7 scores",
        "evidence support rate for interpretability",
        "safety recall and precision for crisis-sensitive language",
        "coverage completeness and multilingual parity",
        "disclosure efficiency: how quickly unresolved items become stable",
    ], font_size=20)

    chart_data = CategoryChartData()
    chart_data.categories = ["Turn 1", "Turn 2", "Turn 3", "Turn 4", "Turn 5", "Turn 6"]
    chart_data.add_series("Adaptive follow-up", (16, 12, 9, 6, 4, 3))
    chart_data.add_series("Fixed script", (16, 14, 12, 10, 8, 6))
    chart = slide.placeholders[10].insert_chart(XL_CHART_TYPE.LINE_MARKERS, chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.maximum_scale = 16
    chart.value_axis.minimum_scale = 0
    chart.category_axis.tick_labels.font.size = Pt(11)
    chart.value_axis.tick_labels.font.size = Pt(11)
    for series, color in zip(chart.series, [UB_BLUE, UB_ORANGE]):
        series.format.line.color.rgb = color
        series.marker.style = XL_MARKER_STYLE.CIRCLE
        series.marker.size = 7
        series.marker.format.fill.solid()
        series.marker.format.fill.fore_color.rgb = color
        series.marker.format.line.color.rgb = color

    add_textbox(slide, Inches(6.65), Inches(5.82), Inches(4.8), Inches(0.28),
                "Illustrative metric behavior only, not experimental results.", font_size=13, color=UB_DK_BLUE, align=PP_ALIGN.CENTER)
    return slide


def add_baselines_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_ONE_COL_BULLETS])
    slide.shapes.title.text = "Baselines and Ablations"
    slide.placeholders[1].text = ""

    add_table(slide, Inches(0.75), Inches(1.95), Inches(12.0), Inches(2.75), [
        ["Comparison", "Adaptive", "Evidence-first", "Safety-aware", "Why include it?"],
        ["Direct questionnaire", "No", "No", "Manual only", "strongest structured baseline"],
        ["Fixed scripted chatbot", "No", "Partial", "Optional", "conversation without confidence tracking"],
        ["Single-pass transcript scoring", "No", "No", "No", "cheap transcript-only scorer"],
        ["No confidence / no safety", "Partial", "Yes", "No", "ablation on design choices"],
    ], header_fill=UB_DK_BLUE)
    add_textbox(slide, Inches(0.9), Inches(5.2), Inches(11.4), Inches(0.55),
                "Main question: does the adaptive evidence-first pipeline improve coverage and interpretability without losing score quality?", font_size=16, color=UB_DK_BLUE)
    return slide


def add_risks_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_ONE_COL_BULLETS])
    slide.shapes.title.text = "Risks, Ethics, and Limitations"
    slide.placeholders[1].text = ""

    # Risk matrix on left
    x0, y0, w, h = Inches(0.7), Inches(2.0), Inches(5.5), Inches(3.4)
    quad_colors = [RGBColor(245, 248, 252), RGBColor(255, 245, 242), RGBColor(238, 247, 248), RGBColor(250, 250, 250)]
    positions = [
        (x0, y0, w/2, h/2, quad_colors[0]),
        (x0 + w/2, y0, w/2, h/2, quad_colors[1]),
        (x0, y0 + h/2, w/2, h/2, quad_colors[2]),
        (x0 + w/2, y0 + h/2, w/2, h/2, quad_colors[3]),
    ]
    for left, top, ww, hh, fill_color in positions:
        rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, ww, hh)
        rect.fill.solid()
        rect.fill.fore_color.rgb = fill_color
        rect.line.color.rgb = UB_LT_GRAY
    add_textbox(slide, Inches(2.4), Inches(1.6), Inches(2.0), Inches(0.25), "Higher impact", font_size=13, color=UB_DK_BLUE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(2.4), Inches(5.45), Inches(2.0), Inches(0.25), "Lower impact", font_size=13, color=UB_DK_BLUE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.25), Inches(3.2), Inches(0.5), Inches(0.4), "Higher\nuncertainty", font_size=12, color=UB_DK_BLUE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(6.0), Inches(3.2), Inches(0.5), Inches(0.4), "Lower\nuncertainty", font_size=12, color=UB_DK_BLUE, align=PP_ALIGN.CENTER)
    risk_labels = [
        ("synthetic realism gap", Inches(1.2), Inches(2.25), UB_ORANGE),
        ("Hindi/code-mix nuance", Inches(3.55), Inches(2.35), UB_BLUE),
        ("label mismatch in public data", Inches(1.0), Inches(4.1), UB_TEAL),
        ("compute limits for 32B model", Inches(3.65), Inches(4.15), UB_LT_BLUE),
    ]
    for label, left, top, color in risk_labels:
        s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, Inches(1.7), Inches(0.45))
        s.fill.solid()
        s.fill.fore_color.rgb = color
        s.line.fill.background()
        add_textbox(slide, left + Inches(0.04), top + Inches(0.06), Inches(1.62), Inches(0.25), label, font_size=11.5, color=WHITE, align=PP_ALIGN.CENTER)

    add_card(slide, Inches(6.8), Inches(2.05), Inches(5.3), Inches(3.35), "Ethical guardrails",
             "Screening support only, not therapy or diagnosis.\nHuman review required for safety-sensitive cases.\nReal-user deployment would require consent, privacy controls, and likely institutional review.\nPhase 1 claims must stay modest because the pilot data is small and largely synthetic.",
             line=UB_ORANGE, title_color=UB_ORANGE)
    return slide


def add_timeline_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_ONE_COL_BULLETS])
    slide.shapes.title.text = "Milestones and Team Split"
    slide.placeholders[1].text = ""

    phases = [
        ("Phase 1", "problem framing\npilot data design\narchitecture + evaluation", UB_BLUE, 0.8),
        ("Phase 2", "prototype pipeline\nconfidence tracker\nsafety classifier", UB_LT_BLUE, 4.3),
        ("Phase 3", "experiments\nablations\nfinal demo + report", UB_ORANGE, 7.8),
    ]
    for title, body, color, left in phases:
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(2.0), Inches(2.9), Inches(1.6))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        add_textbox(slide, Inches(left + 0.15), Inches(2.18), Inches(2.6), Inches(0.3), title, font_size=18, bold=True, color=WHITE)
        add_multiline_box(slide, Inches(left + 0.15), Inches(2.60), Inches(2.6), Inches(0.75), body.split("\n"), font_size=13, color=WHITE)

    for x in [3.75, 7.25]:
        arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(x), Inches(2.55), Inches(0.35), Inches(0.4))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = UB_DK_BLUE
        arrow.line.fill.background()

    add_card(slide, Inches(1.5), Inches(4.4), Inches(4.0), Inches(1.55), "Ritwij",
             "inference engine\nevaluation setup\nintegration and final write-up", line=UB_BLUE)
    add_card(slide, Inches(7.1), Inches(4.4), Inches(4.0), Inches(1.55), "Yash",
             "data design\nannotation workflow\ndialogue logic and presentation", line=UB_LT_BLUE)
    add_textbox(slide, Inches(0.85), Inches(6.12), Inches(11.6), Inches(0.3),
                "Immediate next step after Phase 1: finalize the annotation rubric and build the first pilot batch of conversations.",
                font_size=14.5, color=UB_DK_BLUE, align=PP_ALIGN.CENTER)
    return slide


def add_references_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TWO_COL_BULLETS])
    slide.shapes.title.text = "Selected References"
    set_bullets(slide.placeholders[1], [
        "Kroenke et al. (2001) - PHQ-9",
        "Spitzer et al. (2006) - GAD-7",
        "Dosovitsky et al. (2021) - chatbot PHQ-9",
        "Abd-Alrazaq et al. (2020) - chatbot review",
    ], font_size=20)
    set_bullets(slide.placeholders[13], [
        "Burdisso et al. (2024) - DAIC-WOZ caution",
        "Kakwani et al. (2020) - IndicNLPSuite",
        "Zirikly et al. (2019) - suicide-risk shared task",
    ], font_size=20)
    add_textbox(slide, Inches(0.75), Inches(6.05), Inches(11.3), Inches(0.28),
                "Model/tool sources used in the report: Aya Expanse docs, Mistral NeMo docs, Gemma docs, and LangGraph docs.",
                font_size=14, color=UB_DK_BLUE)
    return slide


def build():
    template_pptx = extract_template()
    prs = Presentation(str(template_pptx))
    clear_all_template_slides(prs)

    add_title_slide(prs)
    add_problem_slide(prs)
    add_gap_questions_slide(prs)
    add_example_slide(prs)
    add_architecture_slide(prs)
    add_data_slide(prs)
    add_annotation_slide(prs)
    add_model_slide(prs)
    add_evaluation_slide(prs)
    add_baselines_slide(prs)
    add_risks_slide(prs)
    add_timeline_slide(prs)
    add_references_slide(prs)

    prs.save(str(OUTPUT))
    print(f"Wrote {OUTPUT}")


if __name__ == "__main__":
    build()
