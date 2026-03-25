from pathlib import Path
from typing import List, Optional, Sequence, Tuple
import sys

from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


DEFAULT_HTML = Path(
    "/Users/ritwij/Library/Containers/net.whatsapp.WhatsApp/Data/tmp/documents/"
    "848E1FB4-3BBB-4F88-8F13-B77D7A9BE0F2/milestone1_slides_v2 (2).html"
)
DEFAULT_OUTPUT = Path("/Users/ritwij/Documents/multilingualChatbot/milestone1_slides_v2.pptx")

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BLUE = RGBColor(0x00, 0x5B, 0xBB)
DARK_BLUE = RGBColor(0x00, 0x3D, 0x8A)
NAVY = RGBColor(0x00, 0x2D, 0x7A)
ACCENT = RGBColor(0x7D, 0xD3, 0xFC)
LIGHT_BLUE_FILL = RGBColor(0xEE, 0xF4, 0xFF)
LIGHT_BLUE_LINE = RGBColor(0xB5, 0xCC, 0xF0)
LIGHT_GREEN_FILL = RGBColor(0xED, 0xFA, 0xF4)
LIGHT_GREEN_LINE = RGBColor(0x8D, 0xD4, 0xB0)
LIGHT_RED_FILL = RGBColor(0xFF, 0xF4, 0xF4)
LIGHT_RED_LINE = RGBColor(0xF0, 0xAA, 0xAA)
LIGHT_GRAY_FILL = RGBColor(0xF5, 0xF8, 0xFD)
LIGHT_GRAY_LINE = RGBColor(0xD5, 0xE2, 0xF5)
TEXT = RGBColor(0x33, 0x33, 0x33)
MUTED = RGBColor(0x66, 0x66, 0x66)
GREEN = RGBColor(0x0A, 0x5E, 0x42)
RED = RGBColor(0xAA, 0x00, 0x00)
ORANGE = RGBColor(0xD0, 0x5A, 0x30)
TEAL = RGBColor(0x1A, 0x7A, 0x9A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_CODE = RGBColor(0x0F, 0x17, 0x2A)
CODE_TEXT = RGBColor(0x94, 0xA3, 0xB8)


def hex_color(value: str) -> RGBColor:
    value = value.strip().lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def ensure_html_looks_right(html_path: Path):
    soup = BeautifulSoup(html_path.read_text(errors="ignore"), "html.parser")
    slide_count = len(soup.select(".slide"))
    if slide_count != 12:
        raise RuntimeError(f"Expected 12 slides in {html_path}, found {slide_count}")


def set_slide_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, shape_type, left, top, width, height, fill: RGBColor, line: Optional[RGBColor] = None, line_width: float = 1.25):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_width)
    return shape


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    font_size: float = 16,
    color: RGBColor = TEXT,
    bold: bool = False,
    italic: bool = False,
    align=PP_ALIGN.LEFT,
    font_name: str = "Arial",
    valign=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return box


def add_lines(
    slide,
    left,
    top,
    width,
    height,
    lines: Sequence[Tuple[str, float, RGBColor, bool, bool]],
    align=PP_ALIGN.LEFT,
    font_name: str = "Arial",
    valign=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = valign
    for idx, (text, size, color, bold, italic) in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name = font_name
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
    return box


def add_header(slide, pill: Optional[str] = None, hero: bool = False):
    bar_fill = NAVY if hero else BLUE
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.58), bar_fill, None)

    circle = add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, Inches(0.28), Inches(0.10), Inches(0.34), Inches(0.34), WHITE, WHITE)
    add_textbox(slide, Inches(0.30), Inches(0.135), Inches(0.30), Inches(0.22), "UB", 11, BLUE, True, False, PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    circle.line.fill.background()

    add_lines(
        slide,
        Inches(0.70),
        Inches(0.11),
        Inches(2.6),
        Inches(0.34),
        [
            ("University at Buffalo", 11, WHITE, True, False),
            ("The State University of New York", 9, WHITE, False, False),
        ],
    )

    if pill:
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(10.7), Inches(0.13), Inches(2.15), Inches(0.24), NAVY if hero else DARK_BLUE, WHITE, 0.75)
        add_textbox(
            slide,
            Inches(10.78),
            Inches(0.145),
            Inches(2.0),
            Inches(0.16),
            pill,
            9,
            WHITE,
            False,
            False,
            PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )

    if not hero:
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, Inches(0.58), SLIDE_W, Inches(0.02), LIGHT_BLUE_LINE, None)


def add_slide_title(slide, title: str):
    add_textbox(slide, Inches(0.52), Inches(0.84), Inches(11.9), Inches(0.34), title, 22, DARK_BLUE, True)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.52), Inches(1.22), Inches(2.2), Inches(0.03), BLUE, None)


def add_note(slide, text: str):
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(6.62), Inches(12.1), Inches(0.015), LIGHT_GRAY_LINE, None)
    add_textbox(slide, Inches(0.55), Inches(6.69), Inches(12.0), Inches(0.35), text, 10.5, BLUE, False, True)


def add_card(
    slide,
    left,
    top,
    width,
    height,
    title: Optional[str],
    body: Optional[str] = None,
    items: Optional[Sequence[str]] = None,
    fill: RGBColor = LIGHT_GRAY_FILL,
    line: RGBColor = LIGHT_GRAY_LINE,
    title_color: RGBColor = DARK_BLUE,
    body_color: RGBColor = TEXT,
):
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height, fill, line)
    body_top = top + Inches(0.18)
    if title:
        add_textbox(slide, left + Inches(0.14), top + Inches(0.11), width - Inches(0.28), Inches(0.18), title, 12.5, title_color, True)
        body_top = top + Inches(0.38)
    if items:
        body_text = "\n".join(f"› {item}" for item in items)
        add_textbox(slide, left + Inches(0.14), body_top, width - Inches(0.28), height - (body_top - top) - Inches(0.10), body_text, 11.5, body_color)
    elif body:
        add_textbox(slide, left + Inches(0.14), body_top, width - Inches(0.28), height - (body_top - top) - Inches(0.10), body, 11.5, body_color)


def add_small_label(slide, left, top, width, text: str, color: RGBColor):
    add_textbox(slide, left, top, width, Inches(0.18), text.upper(), 9.5, color, True)


def add_pill(slide, left, top, width, height, text: str, fill: RGBColor, line: RGBColor, color: RGBColor):
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height, fill, line, 1.0)
    add_textbox(slide, left + Inches(0.05), top + Inches(0.03), width - Inches(0.10), height - Inches(0.06), text, 10.5, color, False, False, PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def add_split_compare(slide, left, top, width, height):
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height, WHITE, LIGHT_GRAY_LINE)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, left + Inches(0.02), top + Inches(0.02), width / 2 - Inches(0.03), height - Inches(0.04), LIGHT_RED_FILL, None)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, left + width / 2 + Inches(0.01), top + Inches(0.02), width / 2 - Inches(0.03), height - Inches(0.04), hex_color("#f0fff8"), None)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, left + width / 2 - Inches(0.005), top + Inches(0.08), Inches(0.01), height - Inches(0.16), LIGHT_GRAY_LINE, None)
    add_small_label(slide, left + Inches(0.14), top + Inches(0.10), Inches(1.2), "Old way", RED)
    add_textbox(slide, left + Inches(0.14), top + Inches(0.30), width / 2 - Inches(0.28), Inches(0.55), '"Have you felt hopeless in the last 2 weeks? (0-3)"', 11.5, TEXT)
    add_small_label(slide, left + width / 2 + Inches(0.12), top + Inches(0.10), Inches(1.4), "ManoVarta", GREEN)
    add_textbox(slide, left + width / 2 + Inches(0.12), top + Inches(0.30), width / 2 - Inches(0.24), Inches(0.55), '"Tell me about how your week has been going..."', 11.5, TEXT)


def add_code_box(slide, left, top, width, height, code: str):
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height, DARK_CODE, DARK_CODE, 0.75)
    add_textbox(slide, left + Inches(0.14), top + Inches(0.12), width - Inches(0.28), height - Inches(0.24), code, 10.3, CODE_TEXT, False, False, PP_ALIGN.LEFT, font_name="Courier New")


def add_table(slide, left, top, width, height, data: Sequence[Sequence[str]], highlight_last: bool = False):
    rows = len(data)
    cols = len(data[0])
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    for c in range(cols):
        table.columns[c].width = int(width / cols)
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c]
            cell.text_frame.word_wrap = True
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = BLUE
            elif highlight_last and r == rows - 1:
                cell.fill.fore_color.rgb = LIGHT_BLUE_FILL
            else:
                cell.fill.fore_color.rgb = WHITE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.name = "Arial"
                    run.font.size = Pt(10.5 if r == 0 else 10.8)
                    run.font.bold = r == 0 or (highlight_last and r == rows - 1 and c == 0)
                    run.font.color.rgb = WHITE if r == 0 else (DARK_BLUE if highlight_last and r == rows - 1 and c == 0 else TEXT)
    return table


def add_arrow_text(slide, left, top, text: str):
    add_textbox(slide, left, top, Inches(0.32), Inches(0.26), text, 19, BLUE, True, False, PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def add_timeline_step(slide, y, title: str, body: str, current: bool = False):
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, Inches(0.75), y + Inches(0.06), Inches(0.12), Inches(0.12), BLUE, None)
    if current:
        ring = add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, Inches(0.71), y + Inches(0.02), Inches(0.20), Inches(0.20), LIGHT_BLUE_FILL, BLUE, 1.0)
        ring.fill.fore_color.rgb = LIGHT_BLUE_FILL
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, Inches(0.75), y + Inches(0.06), Inches(0.12), Inches(0.12), BLUE, None)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.80), y + Inches(0.18), Inches(0.02), Inches(0.58), LIGHT_GRAY_LINE, None)
    add_card(
        slide,
        Inches(0.98),
        y,
        Inches(5.2),
        Inches(0.82),
        title,
        body=body,
        fill=LIGHT_BLUE_FILL if current else LIGHT_GRAY_FILL,
        line=BLUE if current else LIGHT_GRAY_LINE,
    )


def build_title_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BLUE)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(8.8), 0, Inches(4.6), SLIDE_H, DARK_BLUE, None)
    add_header(slide, "Milestone 1 · Phase 1 Proposal", hero=True)
    add_lines(
        slide,
        Inches(0.72),
        Inches(1.75),
        Inches(7.7),
        Inches(2.1),
        [
            ("ManoVarta", 30, WHITE, True, False),
            ("Multilingual Conversational AI", 24, ACCENT, True, False),
            ("for Mental Health Screening", 27, WHITE, True, False),
        ],
    )
    add_textbox(
        slide,
        Inches(0.75),
        Inches(3.78),
        Inches(6.3),
        Inches(0.70),
        "Questionnaire-grounded adaptive dialogue for PHQ-9 / GAD-7 screening in English and Hindi - with evidence-first scoring and parallel safety monitoring.",
        13.5,
        WHITE,
    )
    pills = [
        ("Evidence-first scoring", Inches(0.75), Inches(4.72), Inches(1.95)),
        ("Confidence-based follow-up", Inches(2.82), Inches(4.72), Inches(2.35)),
        ("Parallel safety module", Inches(5.31), Inches(4.72), Inches(1.95)),
        ("English + Hindi + Hinglish", Inches(0.75), Inches(5.22), Inches(2.28)),
        ("LangGraph orchestration", Inches(3.18), Inches(5.22), Inches(2.04)),
    ]
    for text, left, top, width in pills:
        add_pill(slide, left, top, width, Inches(0.34), text, DARK_BLUE, WHITE, WHITE)
    add_lines(
        slide,
        Inches(10.05),
        Inches(6.25),
        Inches(2.45),
        Inches(0.50),
        [
            ("Ritwij & Yash", 11, WHITE, False, False),
            ("[Course Name] · March 2026", 10.5, ACCENT, False, False),
        ],
        align=PP_ALIGN.RIGHT,
    )


def build_problem_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_header(slide)
    add_slide_title(slide, "Why This Problem Matters")
    add_card(slide, Inches(0.60), Inches(1.55), Inches(5.75), Inches(1.00), "Survey fatigue", body="Direct forms feel like a test - patients rush through and give guarded, minimal answers.", fill=LIGHT_RED_FILL, line=LIGHT_RED_LINE, title_color=RED)
    add_card(slide, Inches(0.60), Inches(2.66), Inches(5.75), Inches(1.00), "Social desirability bias", body='Fixed item order misses natural symptom narratives. "Have you felt hopeless?" -> almost always "no."', fill=LIGHT_RED_FILL, line=LIGHT_RED_LINE, title_color=RED)
    add_card(slide, Inches(0.60), Inches(3.77), Inches(5.75), Inches(1.00), "Language gap", body="Most validated tools are English-only. Hindi nuance and code-mixing go unaddressed.", fill=LIGHT_RED_FILL, line=LIGHT_RED_LINE, title_color=RED)

    add_card(
        slide,
        Inches(6.68),
        Inches(1.55),
        Inches(5.95),
        Inches(1.52),
        "What conversation offers",
        items=[
            "Invites more natural disclosure",
            "Follow-up can target ambiguous symptoms",
            "Adaptive coverage without rigid order",
        ],
        fill=LIGHT_GREEN_FILL,
        line=LIGHT_GREEN_LINE,
        title_color=GREEN,
    )
    add_split_compare(slide, Inches(6.68), Inches(3.22), Inches(5.95), Inches(1.56))
    add_note(slide, "Questionnaires remain important. The research question is whether conversational elicitation can preserve PHQ-9/GAD-7 structure while sounding less rigid.")


def build_gap_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_header(slide)
    add_slide_title(slide, "Research Gap and Phase 1 Questions")
    add_small_label(slide, Inches(0.62), Inches(1.50), Inches(2.0), "Known Gaps", RED)
    gap_bodies = [
        "Small pilot data in realistic bilingual screening settings",
        "Public labels rarely map to item-level PHQ-9/GAD-7 scoring",
        "Hindi nuance and code-mixing treated as afterthoughts in existing tools",
    ]
    for idx, body in enumerate(gap_bodies):
        add_card(slide, Inches(0.62), Inches(1.76 + idx * 0.95), Inches(5.65), Inches(0.76), None, body=body, fill=LIGHT_RED_FILL, line=LIGHT_RED_LINE)

    add_small_label(slide, Inches(6.72), Inches(1.50), Inches(3.0), "Phase 1 Questions", DARK_BLUE)
    q_bodies = [
        "Can natural dialogue recover item-level symptom scores with acceptable accuracy?",
        "Does confidence-based follow-up improve coverage vs a fixed scripted chatbot?",
        "Does evidence-first scoring produce more stable outputs than single-pass transcript scoring?",
        "How stable is performance across English, Hindi, and Hinglish?",
    ]
    for idx, body in enumerate(q_bodies):
        add_card(slide, Inches(6.72), Inches(1.76 + idx * 0.78), Inches(5.95), Inches(0.62), None, body=body, fill=LIGHT_BLUE_FILL, line=LIGHT_BLUE_LINE)
    add_note(slide, "Phase 1 goal: a credible research foundation - task definition, data plan, architecture, and evaluation. Not a completed system.")


def build_example_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_header(slide)
    add_slide_title(slide, "Example: Conversation -> Structured Output")
    add_card(
        slide,
        Inches(0.62),
        Inches(1.55),
        Inches(5.35),
        Inches(1.45),
        "Sample patient profile",
        items=[
            "23-year-old graduate student",
            "English-dominant",
            "Moderate sleep disruption + concentration difficulty",
            "No immediate safety flag",
        ],
        fill=LIGHT_BLUE_FILL,
        line=LIGHT_BLUE_LINE,
    )
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.62), Inches(3.16), Inches(5.35), Inches(0.88), LIGHT_GRAY_FILL, LIGHT_GRAY_LINE)
    add_small_label(slide, Inches(0.78), Inches(3.29), Inches(1.3), "User utterance", GREEN)
    add_textbox(slide, Inches(0.78), Inches(3.50), Inches(5.0), Inches(0.34), "\"My sleep schedule is messed up and I can't focus on assignments.\"", 11.5, MUTED, False, True)
    add_card(slide, Inches(0.62), Inches(4.20), Inches(5.35), Inches(1.18), "Next action", body="Ask targeted follow-up about mood or fatigue - sleep and concentration already high-confidence, skip repeating them.", fill=LIGHT_GREEN_FILL, line=LIGHT_GREEN_LINE, title_color=GREEN)
    add_small_label(slide, Inches(6.22), Inches(1.56), Inches(3.0), "Evidence-first structured output", DARK_BLUE)
    code = """{
  "evidence": [
    { "item": "phq_q3_sleep",
      "span": "sleep schedule is messed up",
      "score": 2 },
    { "item": "phq_q7_concentration",
      "span": "can't focus on assignments",
      "score": 2 }
  ],
  "confidence": {
    "phq_q3_sleep": 0.86,
    "phq_q7_concentration": 0.79,
    "phq_q2_low_mood": 0.31  <- unresolved
  },
  "next_action": "follow_up_on_mood"
}"""
    add_code_box(slide, Inches(6.22), Inches(1.82), Inches(6.15), Inches(3.70), code)
    add_note(slide, "Every score is traceable to a quoted evidence span. A correct total with no item evidence is not sufficient for this project.")


def build_architecture_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_header(slide)
    add_slide_title(slide, "Proposed System Architecture")
    box_y = Inches(1.95)
    box_w = Inches(2.65)
    box_h = Inches(1.25)
    specs = [
        (Inches(0.55), LIGHT_GRAY_FILL, LIGHT_GRAY_LINE, "Input", "User Turn", ["English / Hindi", "Hinglish input"], MUTED),
        (Inches(3.45), LIGHT_BLUE_FILL, LIGHT_BLUE_LINE, "Layer 1", "Dialogue Manager", ["Topic steering", "Coverage tracking", "Follow-up strategy"], DARK_BLUE),
        (Inches(6.35), LIGHT_GREEN_FILL, LIGHT_GREEN_LINE, "Layer 2", "Scoring Engine", ["Span extraction", "Item mapping", "0-3 scoring + conf."], GREEN),
        (Inches(9.25), LIGHT_GRAY_FILL, LIGHT_GRAY_LINE, "Output", "Clinician Summary", ["Item scores", "Confidence", "Safety flags"], MUTED),
    ]
    for left, fill, line, label, title, items, title_color in specs:
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, box_y, box_w, box_h, fill, line)
        add_small_label(slide, left + Inches(0.16), box_y + Inches(0.12), Inches(1.2), label, title_color if title_color != MUTED else MUTED)
        add_textbox(slide, left + Inches(0.14), box_y + Inches(0.32), box_w - Inches(0.28), Inches(0.20), title, 12.2, title_color if title_color != MUTED else DARK_BLUE, True)
        add_textbox(slide, left + Inches(0.14), box_y + Inches(0.56), box_w - Inches(0.28), Inches(0.50), "\n".join(f"• {item}" for item in items), 10.8, MUTED)
    add_arrow_text(slide, Inches(3.03), Inches(2.40), "->")
    add_arrow_text(slide, Inches(5.93), Inches(2.40), "->")
    add_arrow_text(slide, Inches(8.83), Inches(2.40), "->")
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.70), Inches(3.62), Inches(1.10), Inches(0.28), LIGHT_RED_FILL, LIGHT_RED_LINE)
    add_textbox(slide, Inches(0.78), Inches(3.68), Inches(0.95), Inches(0.12), "from input v", 10, RED, True, False, PP_ALIGN.CENTER)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1.95), Inches(3.75), Inches(8.40), Inches(0.015), LIGHT_RED_LINE, None)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(10.55), Inches(3.62), Inches(1.20), Inches(0.28), LIGHT_RED_FILL, LIGHT_RED_LINE)
    add_textbox(slide, Inches(10.63), Inches(3.68), Inches(1.04), Inches(0.12), "-> to summary", 10, RED, True, False, PP_ALIGN.CENTER)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.70), Inches(4.18), Inches(11.10), Inches(0.90), LIGHT_RED_FILL, LIGHT_RED_LINE)
    add_lines(
        slide,
        Inches(0.88),
        Inches(4.36),
        Inches(1.40),
        Inches(0.34),
        [
            ("Parallel Safety", 10.5, RED, True, False),
            ("Trigger Module", 10.5, RED, True, False),
        ],
    )
    add_textbox(slide, Inches(2.30), Inches(4.33), Inches(9.10), Inches(0.40), "Runs independently of symptom scoring on every turn. Flags escalation-sensitive language (self-harm, extreme hopelessness) for immediate human review. Does not depend on the scoring engine being correct.", 11.3, MUTED)
    add_note(slide, "Key design choice: safety runs in parallel rather than waiting for the main dialogue model to be correct.")


def build_data_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_header(slide)
    add_slide_title(slide, "Pilot Data Plan")
    add_card(
        slide,
        Inches(0.62),
        Inches(1.55),
        Inches(5.40),
        Inches(1.60),
        "Synthetic conversations - 80 total",
        items=[
            "40 patient profiles x 2 disclosure variants",
            "32 English · 32 Hindi · 16 Hinglish/code-mixed",
            "Profile-guided GPT generation + manual review",
            "Covers minimal -> severe severity range",
        ],
        fill=LIGHT_BLUE_FILL,
        line=LIGHT_BLUE_LINE,
    )
    add_card(
        slide,
        Inches(0.62),
        Inches(3.32),
        Inches(5.40),
        Inches(1.38),
        "Public data (auxiliary only)",
        items=[
            "DAIC-WOZ - few-shot grounding + weak validation",
            "CLPsych / eRisk - safety trigger stress tests",
            "Label mismatch noted: total scores, not item-level",
        ],
        fill=LIGHT_GRAY_FILL,
        line=LIGHT_GRAY_LINE,
    )
    add_small_label(slide, Inches(6.35), Inches(1.58), Inches(3.5), "Dataset composition (80 conversations)", DARK_BLUE)
    bars = [
        ("English", "32 · 40%", BLUE, ACCENT, Inches(1.95)),
        ("Hindi", "32 · 40%", TEAL, ACCENT, Inches(2.52)),
        ("Hinglish / code-mixed", "16 · 20%", ORANGE, hex_color("#ffe0d0"), Inches(3.09)),
    ]
    for label, count, fill, accent, top in bars:
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.35), top, Inches(5.90), Inches(0.42), fill, None)
        add_textbox(slide, Inches(6.55), top + Inches(0.08), Inches(2.90), Inches(0.15), label, 11.5, WHITE, True)
        add_textbox(slide, Inches(10.10), top + Inches(0.06), Inches(1.90), Inches(0.18), count, 13.5, accent, True, False, PP_ALIGN.RIGHT)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.35), Inches(3.85), Inches(5.90), Inches(1.34), LIGHT_GRAY_FILL, LIGHT_GRAY_LINE)
    add_small_label(slide, Inches(6.55), Inches(3.99), Inches(1.6), "Sample annotation", GREEN)
    add_textbox(slide, Inches(6.55), Inches(4.20), Inches(5.40), Inches(0.30), "\"I've been really drained. I used to love the gym but haven't been in weeks.\"", 11.2, MUTED, False, True)
    chips = [
        ("Anhedonia -> 2", Inches(6.55), Inches(4.70), Inches(1.40)),
        ("Fatigue -> 3", Inches(8.12), Inches(4.70), Inches(1.22)),
        ("PHQ-9 ~= 9/27", Inches(9.50), Inches(4.70), Inches(1.55)),
    ]
    for text, left, top, width in chips:
        add_pill(slide, left, top, width, Inches(0.28), text, LIGHT_BLUE_FILL, LIGHT_BLUE_LINE, DARK_BLUE)
    add_note(slide, "Planned pilot composition only - not presented as a completed clinical dataset.")


def build_annotation_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_header(slide)
    add_slide_title(slide, "Annotation Protocol and Data Schema")
    top_cards = [
        ("1. Profile design", "Age, context, symptom pattern, disclosure style", LIGHT_GRAY_FILL, LIGHT_GRAY_LINE),
        ("2. Dialogue drafting", "Guided generation followed by manual editing", LIGHT_GRAY_FILL, LIGHT_GRAY_LINE),
        ("3. Double annotation", "PHQ-9, GAD-7, evidence spans, safety flag - independently", LIGHT_BLUE_FILL, LIGHT_BLUE_LINE),
        ("4. Consensus pass", "Resolve disagreements against DSM-5-TR, freeze gold labels", LIGHT_BLUE_FILL, LIGHT_BLUE_LINE),
    ]
    for idx, (title, body, fill, line) in enumerate(top_cards):
        add_card(slide, Inches(0.62 + idx * 3.12), Inches(1.55), Inches(2.86), Inches(1.02), title, body=body, fill=fill, line=line)
    add_small_label(slide, Inches(0.62), Inches(2.88), Inches(2.0), "Gold label schema", DARK_BLUE)
    code = """{
  "patient_id": "MB-P001",
  "language": "en",
  "evidence_spans": ["sleep schedule is messed up"],
  "phq_q3_sleep": 2,
  "phq_q7_concentration": 2,
  "safety_flag": "none"
}"""
    add_code_box(slide, Inches(0.62), Inches(3.12), Inches(12.00), Inches(2.10), code)
    add_note(slide, "Agreement targets: weighted Cohen's Kappa for item labels · span overlap for evidence · explicit notes on ambiguity or contradiction")


def build_model_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_header(slide)
    add_slide_title(slide, "Model Stack and Compute Rationale")
    model_specs = [
        ("Primary", "Dialogue + Evidence", "Aya Expanse 32B", LIGHT_BLUE_FILL, BLUE, DARK_BLUE),
        ("Comparison", "Open Baseline", "Mistral NeMo 12B", LIGHT_GRAY_FILL, hex_color("#7aace0"), DARK_BLUE),
        ("Optional", "Second Baseline", "Gemma 3 12B", LIGHT_GRAY_FILL, hex_color("#9ac0e0"), DARK_BLUE),
        ("Utility", "Safety / Retrieval", "IndicBERT-style encoder", LIGHT_RED_FILL, hex_color("#e09090"), RED),
    ]
    for idx, (label, title, body, fill, line, title_color) in enumerate(model_specs):
        left = Inches(0.62 + idx * 3.06)
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, Inches(1.55), Inches(2.84), Inches(1.12), fill, line, 1.75)
        add_pill(slide, left + Inches(0.12), Inches(1.68), Inches(0.95), Inches(0.22), label, WHITE if label == "Utility" else LIGHT_BLUE_FILL, line, title_color)
        add_textbox(slide, left + Inches(0.12), Inches(1.98), Inches(2.55), Inches(0.20), title, 12.5, title_color, True)
        add_textbox(slide, left + Inches(0.12), Inches(2.24), Inches(2.55), Inches(0.18), body, 11.5, title_color)
    add_card(slide, Inches(0.62), Inches(3.05), Inches(5.90), Inches(1.35), "Why modular?", body="Dialogue, scoring, and safety are fully separable - each layer can be swapped, inspected, and ablated independently without touching the others.", fill=LIGHT_BLUE_FILL, line=LIGHT_BLUE_LINE)
    add_card(slide, Inches(6.72), Inches(3.05), Inches(5.90), Inches(1.35), "Compute plan", items=["Smaller models (Mistral, Gemma) handle most iteration", "Aya Expanse for targeted evaluation subsets via quantized / hosted inference", "No full supervised fine-tuning assumed in Phase 1"], fill=LIGHT_GRAY_FILL, line=LIGHT_GRAY_LINE)


def build_evaluation_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_header(slide)
    add_slide_title(slide, "Evaluation Plan")
    left_cards = [
        ("Core scoring accuracy", "Item-level MAE and macro-F1 for PHQ-9/GAD-7 scores vs gold labels", LIGHT_BLUE_FILL, LIGHT_BLUE_LINE, DARK_BLUE),
        ("Interpretability & safety", "Evidence support rate · safety recall and precision for crisis language", LIGHT_GREEN_FILL, LIGHT_GREEN_LINE, GREEN),
        ("Linguistic equity", "Coverage completeness and multilingual parity across English / Hindi / Hinglish", LIGHT_GRAY_FILL, LIGHT_GRAY_LINE, DARK_BLUE),
        ("Operational efficiency", "Disclosure efficiency - turns for unresolved items to reach stable confidence", LIGHT_GRAY_FILL, LIGHT_GRAY_LINE, DARK_BLUE),
    ]
    for idx, (title, body, fill, line, title_color) in enumerate(left_cards):
        add_card(slide, Inches(0.62), Inches(1.58 + idx * 0.90), Inches(5.30), Inches(0.76), title, body=body, fill=fill, line=line, title_color=title_color)
    add_small_label(slide, Inches(6.25), Inches(1.58), Inches(2.5), "Baselines & Ablations", DARK_BLUE)
    data = [
        ["Comparison", "Adaptive", "Evidence-first", "Safety"],
        ["Direct questionnaire", "No", "No", "Manual"],
        ["Fixed scripted chatbot", "No", "Partial", "Optional"],
        ["Single-pass scoring", "No", "No", "No"],
        ["No conf / no safety", "Partial", "Yes", "No"],
        ["ManoVarta (ours)", "Yes", "Yes", "Yes"],
    ]
    add_table(slide, Inches(6.25), Inches(1.86), Inches(6.05), Inches(2.58), data, highlight_last=True)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(6.25), Inches(4.60), Inches(6.05), Inches(0.015), LIGHT_GRAY_LINE, None)
    add_textbox(slide, Inches(6.25), Inches(4.70), Inches(6.00), Inches(0.32), "Main question: does the adaptive evidence-first pipeline improve coverage and interpretability without losing score quality?", 10.8, BLUE, False, True)


def build_risk_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_header(slide)
    add_slide_title(slide, "Risks, Ethics, and Limitations")
    q_left = Inches(0.70)
    q_top = Inches(1.75)
    q_w = Inches(5.60)
    q_h = Inches(3.00)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, q_left, q_top, q_w, q_h, hex_color("#fafbfd"), LIGHT_GRAY_LINE)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, q_left + q_w / 2, q_top + Inches(0.12), Inches(0.02), q_h - Inches(0.24), LIGHT_GRAY_LINE, None)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, q_left + Inches(0.12), q_top + q_h / 2, q_w - Inches(0.24), Inches(0.02), LIGHT_GRAY_LINE, None)
    add_textbox(slide, Inches(2.45), Inches(1.53), Inches(2.0), Inches(0.18), "Higher Impact", 9.5, MUTED, True, False, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(2.45), Inches(4.78), Inches(2.0), Inches(0.18), "Lower Impact", 9.5, MUTED, True, False, PP_ALIGN.CENTER)
    left_label = add_textbox(slide, Inches(0.36), Inches(3.02), Inches(0.60), Inches(0.24), "High Unc.", 9.5, MUTED, True, False, PP_ALIGN.CENTER)
    left_label.rotation = 270
    right_label = add_textbox(slide, Inches(5.92), Inches(3.02), Inches(0.60), Inches(0.24), "Low Unc.", 9.5, MUTED, True, False, PP_ALIGN.CENTER)
    right_label.rotation = 90
    risks = [
        ("Synthetic\nRealism Gap", Inches(1.05), Inches(2.05), LIGHT_RED_LINE, RED),
        ("Hindi / Code-mix\nNuance", Inches(4.00), Inches(2.05), hex_color("#aac8f0"), NAVY),
        ("Label Mismatch\nin Public Data", Inches(1.05), Inches(3.65), hex_color("#9ad4b8"), GREEN),
        ("Compute Limits\n(32B Model)", Inches(4.00), Inches(3.65), hex_color("#b8d8e8"), hex_color("#1a4a5a")),
    ]
    for text, left, top, fill, color in risks:
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, Inches(1.72), Inches(0.58), fill, None)
        add_textbox(slide, left + Inches(0.05), top + Inches(0.08), Inches(1.62), Inches(0.36), text, 10.0, color, True, False, PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_small_label(slide, Inches(6.85), Inches(1.74), Inches(2.3), "Ethical Guardrails", RED)
    ethical = [
        "Screening support only - not for therapy or diagnosis.",
        "Human review required for all safety-sensitive cases.",
        "Deployment requires consent, privacy controls, and IRB.",
        "Phase 1 claims remain modest due to small synthetic pilot.",
    ]
    for idx, text in enumerate(ethical):
        add_card(slide, Inches(6.85), Inches(2.05 + idx * 0.73), Inches(5.40), Inches(0.58), None, body=text, fill=LIGHT_RED_FILL, line=LIGHT_RED_LINE)


def build_milestone_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_header(slide)
    add_slide_title(slide, "Milestones and Team Split")
    add_timeline_step(slide, Inches(1.72), "Phase 1 - Proposal (now)", "Problem framing · pilot data design · architecture + evaluation plan · this presentation", True)
    add_timeline_step(slide, Inches(2.72), "Phase 2 - Prototype", "Inference pipeline · confidence tracker · safety classifier · first end-to-end runs", False)
    add_timeline_step(slide, Inches(3.72), "Phase 3 - Experiments & Demo", "Baseline comparisons · ablations · error analysis · final report + demo", False)
    add_card(slide, Inches(6.80), Inches(1.78), Inches(5.30), Inches(1.15), "Ritwij", items=["Inference engine and scoring logic", "Evaluation setup and metrics", "System integration and final write-up"], fill=LIGHT_BLUE_FILL, line=LIGHT_BLUE_LINE)
    add_card(slide, Inches(6.80), Inches(3.08), Inches(5.30), Inches(1.15), "Yash", items=["Data design and annotation workflow", "Dialogue logic and architecture framing", "Presentation preparation"], fill=LIGHT_GREEN_FILL, line=LIGHT_GREEN_LINE, title_color=GREEN)
    add_textbox(slide, Inches(6.82), Inches(4.56), Inches(5.25), Inches(0.18), "Shared: rubric design · consensus annotation · ethics review · demo rehearsal", 10.5, BLUE, False, True)
    add_textbox(slide, Inches(6.82), Inches(4.88), Inches(5.25), Inches(0.24), "Immediate next step: finalize annotation rubric and build first pilot batch.", 10.5, BLUE, False, True)


def build_thanks_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BLUE)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(8.8), 0, Inches(4.6), SLIDE_H, DARK_BLUE, None)
    add_header(slide, hero=True)
    add_textbox(slide, Inches(4.20), Inches(2.35), Inches(4.8), Inches(0.60), "Thank You", 34, WHITE, True, False, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(3.00), Inches(3.05), Inches(7.2), Inches(0.48), "ManoVarta - keeping clinical structure inside a human conversation.", 15, WHITE, False, False, PP_ALIGN.CENTER)
    add_pill(slide, Inches(4.08), Inches(4.05), Inches(1.20), Inches(0.34), "Ritwij", DARK_BLUE, WHITE, WHITE)
    add_pill(slide, Inches(5.55), Inches(4.05), Inches(1.10), Inches(0.34), "Yash", DARK_BLUE, WHITE, WHITE)
    add_pill(slide, Inches(6.92), Inches(4.05), Inches(2.40), Inches(0.34), "[Course Name] · March 2026", DARK_BLUE, WHITE, WHITE)


def build_deck(output_path: Path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_title_slide(prs)
    build_problem_slide(prs)
    build_gap_slide(prs)
    build_example_slide(prs)
    build_architecture_slide(prs)
    build_data_slide(prs)
    build_annotation_slide(prs)
    build_model_slide(prs)
    build_evaluation_slide(prs)
    build_risk_slide(prs)
    build_milestone_slide(prs)
    build_thanks_slide(prs)

    prs.save(str(output_path))


def main():
    html_path = Path(sys.argv[1]) if len(sys.argv) > 1 else DEFAULT_HTML
    output_path = Path(sys.argv[2]) if len(sys.argv) > 2 else DEFAULT_OUTPUT
    ensure_html_looks_right(html_path)
    build_deck(output_path)
    print(f"Wrote editable PPT to {output_path}")


if __name__ == "__main__":
    main()
